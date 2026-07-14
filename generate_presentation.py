import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color scheme
    CORAL_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
    AMBER_GOLD = RGBColor(0xFF, 0x9F, 0x1C)
    SLATE_BLACK = RGBColor(0x21, 0x25, 0x29)
    WARM_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    def apply_title_style(title_shape, text):
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = 'Helvetica'
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = CORAL_ORANGE
            
    def apply_content_style(content_shape, bullet_points):
        tf = content_shape.text_frame
        tf.clear()
        for idx, point in enumerate(bullet_points):
            p = tf.add_paragraph()
            p.text = point
            p.font.name = 'Helvetica'
            p.font.size = Pt(16)
            p.font.color.rgb = SLATE_BLACK
            p.space_after = Pt(12)
            if idx > 0:
                p.level = 0

    # 1. Slide: Title
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "FoodieGo 🍔"
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = 'Helvetica'
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = CORAL_ORANGE
        paragraph.alignment = PP_ALIGN.CENTER
        
    subtitle.text = "Task 5: Finalization, Deployment & Presentation\nAndroid Development Internship Project"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.name = 'Helvetica'
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = SLATE_BLACK
        paragraph.alignment = PP_ALIGN.CENTER

    # 2. Slide: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Project Introduction")
    apply_content_style(slide.shapes.placeholders[1], [
        "FoodieGo is a high-performance, premium food ordering Android application.",
        "Engineered with a hybrid architecture bridging a local Node.js API with Firebase backend services.",
        "Adheres strictly to Google Material Design 3 guidelines for fluid transitions, outlines, and micro-animations.",
        "Designed to offer Zomato/Swiggy-grade user experience with strong offline resilience."
    ])

    # 3. Slide: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Problem Statement")
    apply_content_style(slide.shapes.placeholders[1], [
        "Traditional food apps display slow loading states and crash completely during network fluctuations.",
        "Inconsistent layout designs distract users, leading to high cart abandonment rates.",
        "Lack of local caching isolates users from viewing past orders or browsing menus while offline.",
        "Navigation flow bugs (e.g. notification redirects, backstack loopings) cause frustration."
    ])

    # 4. Slide: Project Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Project Objectives")
    apply_content_style(slide.shapes.placeholders[1], [
        "Deliver a highly consistent Material Design 3 interface with outline boxes, loaders, and transitions.",
        "Implement a hybrid synchronization repository: local Node.js databases synced with Firebase Firestore.",
        "Build 100% offline cache capabilities for menus, active carts, and order histories.",
        "Prepare Google Play Store submission deliverables including signed release APK and design graphics."
    ])

    # 5. Slide: Premium Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Key Application Features")
    apply_content_style(slide.shapes.placeholders[1], [
        "Splash Screen: Diagnostic 2-second diagonal gradient transition with brand cloche animation.",
        "Filters & Sorting: Real-time search by food, category, rating, veg, fast delivery, and price ranks.",
        "Interactive Cart & Checkout: Live bill detail sums (subtotal, taxes, delivery) and address selections.",
        "Timeline Order Tracker: Visual multi-step progress bar (Placed -> Preparing -> Out -> Delivered).",
        "In-App Notification Log: Dedicated bottom sheet alert history and localized system notifications."
    ])

    # 6. Slide: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Technology Stack")
    apply_content_style(slide.shapes.placeholders[1], [
        "Frontend Engine: Java (JDK 8 / 1.8 compatibility) and XML layouts (ConstraintLayout).",
        "Networking Library: Retrofit 2 for async REST API calls with Gson deserialization.",
        "Media Processing: Glide v4 for Unsplash photo caching and scaling.",
        "Animations Library: Lottie Android for rich vector JSON loading animations.",
        "Firebase Infrastructure: Firebase Auth, Cloud Firestore, Firebase Storage, and FCM Messaging."
    ])

    # 7. Slide: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "System Architecture")
    apply_content_style(slide.shapes.placeholders[1], [
        "Modular Pattern: Organized cleanly by packages (activities, adapters, models, network, services, utils).",
        "View Binding: 100% type-safe view binding variables in XML lookups preventing NullPointerExceptions.",
        "Repository Pattern: Mediates database calls; queries Node.js first, falling back to cached JSON or Firestore.",
        "Tactile Feedback: Custom programmatic scale triggers on major buttons for tactile response."
    ])

    # 8. Slide: Firebase Integration Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Firebase Integrations")
    apply_content_style(slide.shapes.placeholders[1], [
        "Firebase Auth: Secures user credentials and maps UIDs across platforms.",
        "Cloud Firestore: Stores user metadata, custom address books, and order entries.",
        "Firebase Storage: Hosts high-resolution profile pictures (profile_images/{uid}.jpg).",
        "Firebase Cloud Messaging (FCM): Dispatches push alerts triggered on order milestones."
    ])

    # 9. Slide: Node.js API Integration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Node.js REST API Integration")
    apply_content_style(slide.shapes.placeholders[1], [
        "Backend Server: Custom Express server running on local port 8001.",
        "Database Layer: MongoDB (mongoose) storing catalog foods, carts, and order entries.",
        "Automatic Seeding: Automatically seeds default foods on database connections.",
        "Endpoints mapping: /api/products, /api/auth/login, /api/cart, /api/orders."
    ])

    # 10. Slide: Challenges Faced
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Challenges Faced")
    apply_content_style(slide.shapes.placeholders[1], [
        "Lottie Asset Crash: Animation files inside res/raw caused runtime crashes when loaded via string path.",
        "Navigation Misalignments: Notification bell opened Favorites; no central UI for notification logs existed.",
        "Locale Bug: Switching app to Hindi resulted in English text due to missing resources folder.",
        "Release Signing: Standard debug builds are not production-ready; required signed configurations."
    ])

    # 11. Slide: Solutions Implemented
    slide = prs.shapes = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Solutions Implemented")
    apply_content_style(slide.shapes.placeholders[1], [
        "Assets Folder Setup: Created assets/ directory and duplicated JSON files allowing string load by name.",
        "Notifications Bottom Sheet: Bound bell to a custom BottomSheetDialogFragment rendering logged alerts.",
        "Favorites Card & Bindings: Inserted a 'My Favorites' Card in Profile settings and mapped it to FavoritesActivity.",
        "Hindi Localization: Added values-hi/strings.xml to render Hindi texts when locale changes.",
        "Keystore Generation: Created release.keystore and integrated signing configuration in build.gradle."
    ])

    # 12. Slide: Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Future Scope")
    apply_content_style(slide.shapes.placeholders[1], [
        "Real-Time Map Tracking: Integrate Google Maps SDK to trace delivery riders.",
        "Payment Gateway: Add Razorpay / Stripe SDKs for live transactions.",
        "Loyalty points: Reward points system for frequent order completions.",
        "AI Recommendation: Gemini API recommendations of popular foods based on history."
    ])

    # 13. Slide: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_style(slide.shapes.title, "Conclusion")
    apply_content_style(slide.shapes.placeholders[1], [
        "FoodieGo has been successfully refactored and finalized into a production-grade app.",
        "Compiles with 0 warnings, zero memory crashes, and full support for Dark Mode/Localization.",
        "Release bundle configurations are set up and fully prepared for Google Play Store upload.",
        "Deliverables are complete and ready for submission."
    ])

    # Save presentation
    output_path = "presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully as {output_path}!")

if __name__ == '__main__':
    create_presentation()
